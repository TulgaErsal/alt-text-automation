#!/usr/bin/env python3
"""
generate_alt_text.py

Batch-processes all pictures in a PowerPoint file using the UMich AI Helper
Alt Text Generator tool, authenticating via institutional Google sign-in.

The script opens Firefox, navigates to the tool URL, and waits for you to
complete Google authentication. After sign-in it automates form submission
for each image and writes the generated alt text back into the .pptx file.

Requirements:
    pip install selenium python-pptx Pillow

Usage:
    python generate_alt_text.py slides.pptx --url https://your-tool-url/alt-text-generator
    python generate_alt_text.py slides.pptx --url https://... --version long
    python generate_alt_text.py slides.pptx --url https://... --output slides_accessible.pptx
"""

import argparse
import io
import re
import sys
import tempfile
import threading
import time
from pathlib import Path

# ── dependency checks ─────────────────────────────────────────────────────────

try:
    from selenium import webdriver
    from selenium.common.exceptions import (
        InvalidSessionIdException,
        NoSuchWindowException,
        SessionNotCreatedException,
        TimeoutException,
    )
    # Exceptions that mean the browser process is gone — not a per-image error
    _SESSION_LOST_EXCEPTIONS = (InvalidSessionIdException, NoSuchWindowException)
    from selenium.webdriver.chrome.options import Options as ChromeOptions
    from selenium.webdriver.chrome.service import Service as ChromeService
    from selenium.webdriver.common.by import By
    from selenium.webdriver.edge.options import Options as EdgeOptions
    from selenium.webdriver.edge.service import Service as EdgeService
    from selenium.webdriver.firefox.options import Options as FirefoxOptions
    from selenium.webdriver.firefox.service import Service as FirefoxService
    from selenium.webdriver.support import expected_conditions as EC
    from selenium.webdriver.support.ui import WebDriverWait
except ImportError:
    sys.exit("Missing dependency: pip install selenium")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("Missing dependency: pip install python-pptx")

try:
    from PIL import Image
except ImportError:
    sys.exit("Missing dependency: pip install Pillow")

# ── constants ─────────────────────────────────────────────────────────────────

SCRIPT_DIR = Path(__file__).parent

# geckodriver bundled in the project folder; fall back to system PATH
_LOCAL_GECKODRIVER = SCRIPT_DIR / "geckodriver.exe"

# Formats the web form accepts
ACCEPTED_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}

# PPTX content-types that need PIL conversion before the web form will accept them
NEEDS_CONVERSION = {"image/bmp", "image/tiff", "image/x-bmp"}

# Vector formats that require special conversion before upload
VECTOR_TYPES = {"image/x-emf", "image/x-wmf", "image/svg+xml"}

# How long to wait for authentication (seconds)
AUTH_TIMEOUT = 300

# How long to wait for AI generation to complete (seconds)
GENERATION_TIMEOUT = 60

# Seconds of text-stability required before we consider streaming done
STABILITY_SECONDS = 1.0

# ── web form option maps ──────────────────────────────────────────────────────
# Keys are the short CLI names; values are the aria-label / search strings
# used to interact with the Vuetify form controls.

PURPOSE_OPTIONS: dict[str, str] = {
    "general":       "General Description (for decorative or illustrative images)",
    "educational":   "Detailed Educational Description (for charts, diagrams, maps etc)",
    "instructional": "Instructional Use (used in tutorials or learning modules)",
    "marketing":     "Marketing or Promotional Use",
    "icon":          "Interface Icon or Button",
}

INCLUDE_OPTIONS: dict[str, str] = {
    "data-values": "Include Data Values",
    "captions":    "Include Captions/Labels",
}

# Distinctive substring to match the dropdown item text for each tone
TONE_OPTIONS: dict[str, str] = {
    "formal":         "formal tone",
    "academic":       "academic style",
    "professional":   "professional tone",
    "neutral":        "neutral tone",
    "conversational": "conversational tone",
    "casual":         "casual tone",
    "colloquial":     "colloquial tone",
}


# ── image helpers ─────────────────────────────────────────────────────────────

def convert_to_png(image_bytes: bytes) -> bytes:
    """Re-encode any PIL-readable raster image as PNG."""
    with Image.open(io.BytesIO(image_bytes)) as img:
        buf = io.BytesIO()
        img.convert("RGB").save(buf, format="PNG")
        return buf.getvalue()


def _convert_svg_to_png(image_bytes: bytes) -> bytes:
    """Convert SVG to PNG using cairosvg (pip install cairosvg)."""
    try:
        import cairosvg
    except ImportError:
        raise RuntimeError(
            "SVG conversion requires cairosvg: pip install cairosvg"
        )
    return cairosvg.svg2png(bytestring=image_bytes)


def _convert_metafile_to_png(image_bytes: bytes, is_emf: bool, _debug: bool = False) -> bytes:
    """
    Render an EMF or WMF to PNG using PIL.
    """
    with Image.open(io.BytesIO(image_bytes)) as img:
        out = io.BytesIO()
        img.save(out, format="PNG")
        return out.getvalue()


def save_image_to_temp(image_bytes: bytes, content_type: str) -> Path:
    """
    Write image bytes to a named temp file and return its path.
    Converts to PNG when the format is not accepted by the web form.
    """
    ct = content_type.lower()

    if ct == "image/svg+xml":
        image_bytes = _convert_svg_to_png(image_bytes)
        suffix = ".png"
    elif ct in {"image/x-emf", "image/x-wmf"}:
        # Detect the actual format from magic bytes rather than trusting the
        # content type — some pptx files mislabel EMF files as image/x-wmf.
        # EMF starts with EMR_HEADER (iType=1): bytes 01 00 00 00.
        # APM/WMF starts with 0x9AC6CDD7.
        import struct as _struct
        emf_magic = len(image_bytes) >= 4 and _struct.unpack_from("<I", image_bytes)[0] == 0x00000001
        image_bytes = _convert_metafile_to_png(image_bytes, is_emf=emf_magic)
        suffix = ".png"
    elif ct in NEEDS_CONVERSION or ct not in {
        "image/jpeg", "image/jpg", "image/png", "image/gif", "image/webp"
    }:
        image_bytes = convert_to_png(image_bytes)
        suffix = ".png"
    else:
        ext_map = {
            "image/jpeg": ".jpg",
            "image/jpg": ".jpg",
            "image/png": ".png",
            "image/gif": ".gif",
            "image/webp": ".webp",
        }
        suffix = ext_map.get(ct, ".png")

    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    tmp.write(image_bytes)
    tmp.close()
    return Path(tmp.name)


# Heading pattern shared by the streaming detector and the response extractor.
# Matches a standalone line that starts with Short/Medium/Long (after optional
# bold markers) and contains the word "version" anywhere on the same line.
# This handles the standard format ("**Long Version**:") as well as style
# annotations the AI sometimes adds ("Long (academic) version:",
# "Long version (academic style):").
# Deliberately does NOT match inline phrases like
# "a short version, medium version, and long version".
_VERSION_HEADING = re.compile(
    r"(?im)^\s*[*]{0,3}\s*(Short|Medium|Long)\b[^\n]*\b(?:version|alt\s+text)\b[^\n]*$"
)


def _extract_response_text(page_text: str) -> str | None:
    """
    Isolate the AI response block from a full page innerText dump.

    The AI always outputs the three versions in Short → Medium → Long order.
    Find the first occurrence of that exact sequence and return everything
    from the 'Short Version' heading onwards.  Returns None when the sequence
    has not appeared yet (AI still streaming), so callers can keep waiting.
    """
    headings = list(_VERSION_HEADING.finditer(page_text))
    for i in range(len(headings) - 2):
        a = headings[i].group(1).lower()
        b = headings[i + 1].group(1).lower()
        c = headings[i + 2].group(1).lower()
        if (a, b, c) == ("short", "medium", "long"):
            return page_text[headings[i].start():]
    return None


# ── alt text parsing ──────────────────────────────────────────────────────────

def extract_version(full_text: str, version: str) -> str:
    """
    Pull the short / medium / long section out of the three-version response.
    The tool always outputs headers like '**Short Version**' or 'Short Version:'.
    Falls back to a 500-char truncation of the whole response.
    """
    heading = version.capitalize()
    # Anchor to the start of a line ((?m) flag) so we never match "long" or
    # "short" appearing inline in preamble sentences like "long description:".
    # Format matches _VERSION_HEADING used by the streaming detector.
    pattern = (
        rf"(?im)^\s*[*]{{0,3}}\s*{heading}\b[^\n]*\b(?:version|alt\s+text)\b[^\n]*$"
        rf"\n*"
        rf"(.*?)"
        rf"(?=\n\s*[*]{{0,3}}\s*(?:Short|Medium|Long)\b[^\n]*\b(?:version|alt\s+text)\b|\Z)"
    )
    match = re.search(pattern, full_text, re.DOTALL)
    if match:
        text = match.group(1).strip()
        text = re.sub(r"\*{1,3}(.*?)\*{1,3}", r"\1", text)   # strip markdown
        first_para = text.split("\n\n")[0].strip()
        if first_para:
            return first_para

    # Fallback
    plain = re.sub(r"\*{1,3}(.*?)\*{1,3}", r"\1", full_text).strip()
    return plain[:500]


# ── PPTX alt text writer ──────────────────────────────────────────────────────

def set_alt_text(shape, description: str, title: str = "") -> bool:
    """
    Write alt text directly into the shape's XML.

    PowerPoint stores alt text on <p:cNvPr>:
      'title' → Title field shown in the alt text dialog
      'descr' → Description field read by screen readers
    """
    cNvPr = None
    for attr in ("nvPicPr", "nvSpPr", "nvGrpSpPr"):
        try:
            cNvPr = getattr(shape._element, attr).cNvPr
            break
        except AttributeError:
            continue

    if cNvPr is None:
        return False

    cNvPr.set("title", title[:255])
    cNvPr.set("descr", description[:2000])
    return True


# ── browser automation ────────────────────────────────────────────────────────

# Common Firefox installation paths on Windows
_FIREFOX_CANDIDATES = [
    r"C:\Program Files\Mozilla Firefox\firefox.exe",
    r"C:\Program Files (x86)\Mozilla Firefox\firefox.exe",
]


def _build_firefox(geckodriver_path: Path | None) -> webdriver.Firefox:
    options = FirefoxOptions()
    options.headless = False

    # Locate the Firefox binary if it is not on the system PATH
    for candidate in _FIREFOX_CANDIDATES:
        if Path(candidate).exists():
            options.binary_location = candidate
            break

    # Suppress form-data restore dialogs.
    options.set_preference("browser.formfill.enable", False)
    options.set_preference("browser.sessionstore.resume_from_crash", False)
    # Silently dismiss any unexpected confirm/alert dialogs (e.g. "load saved
    # form data?") instead of raising UnexpectedAlertPresentException.
    options.set_capability("unhandledPromptBehavior", "dismiss")

    service = (
        FirefoxService(executable_path=str(geckodriver_path))
        if geckodriver_path and geckodriver_path.exists()
        else FirefoxService()
    )
    return webdriver.Firefox(service=service, options=options)


def _build_chrome() -> webdriver.Chrome:
    options = ChromeOptions()
    options.add_argument("--start-maximized")
    options.set_capability("unhandledPromptBehavior", "dismiss")
    # Selenium Manager (bundled with Selenium 4.6+) downloads chromedriver automatically
    return webdriver.Chrome(service=ChromeService(), options=options)


def _build_edge() -> webdriver.Edge:
    options = EdgeOptions()
    options.add_argument("--start-maximized")
    options.set_capability("unhandledPromptBehavior", "dismiss")
    # Selenium Manager downloads msedgedriver automatically
    return webdriver.Edge(service=EdgeService(), options=options)


def build_driver(browser: str, geckodriver_path: Path | None) -> webdriver.Remote:
    """
    Launch the requested browser. 'auto' tries Firefox → Chrome → Edge.
    Raises SystemExit with a helpful message if no browser can be started.
    """
    builders = {
        "firefox": lambda: _build_firefox(geckodriver_path),
        "chrome":  _build_chrome,
        "edge":    _build_edge,
    }

    if browser != "auto":
        try:
            drv = builders[browser]()
            print(f"Browser  : {browser}")
            return drv
        except SessionNotCreatedException as exc:
            sys.exit(
                f"Could not launch {browser}: {exc.msg}\n"
                "Check that the browser is installed, or choose a different one "
                "with --browser."
            )

    # auto: try each in order
    for name, builder in builders.items():
        try:
            drv = builder()
            print(f"Browser  : {name} (auto-detected)")
            return drv
        except (SessionNotCreatedException, Exception):
            continue

    sys.exit(
        "Could not launch any browser (tried Firefox, Chrome, Edge).\n"
        "Install one of those browsers and re-run, or specify --browser explicitly."
    )


def wait_for_auth(driver: webdriver.Firefox, url: str, raise_on_timeout: bool = False) -> None:
    """
    Navigate to the tool and block until the user has authenticated.
    Authentication is confirmed when the profile link appears in the header.

    When raise_on_timeout is True, raises TimeoutError instead of calling
    sys.exit (used by the GUI so it can handle the failure gracefully).
    """
    print(f"\nOpening browser → {url}")
    driver.get(url)

    sep = "=" * 60
    print(f"\n{sep}\nPlease sign in with your institutional Google account.\nWaiting up to {AUTH_TIMEOUT // 60} minutes …\n{sep}")

    try:
        # The profile link (e.g. 'Tulga Ersal') appears in .top-user once logged in
        WebDriverWait(driver, AUTH_TIMEOUT).until(
            EC.presence_of_element_located((By.CSS_SELECTOR, 'a[href="/profile"]'))
        )
    except TimeoutException:
        if raise_on_timeout:
            raise TimeoutError("Timed out waiting for authentication.")
        driver.quit()
        sys.exit("Timed out waiting for authentication. Please re-run the script.")

    print("Authentication confirmed. Starting batch processing …\n")


def reset_page(driver: webdriver.Remote, url: str) -> None:
    """Navigate back to the tool page to get a fresh form."""
    driver.get(url)
    # Dismiss any residual browser prompt (e.g. "load saved form data?")
    try:
        driver.switch_to.alert.dismiss()
    except Exception:
        pass
    # Wait for the file input to be present before proceeding
    WebDriverWait(driver, 30).until(
        EC.presence_of_element_located((By.CSS_SELECTOR, 'input[type="file"]'))
    )


def _set_form_options(
    driver: webdriver.Remote,
    wait: WebDriverWait,
    purpose: str | None,
    includes: list[str],
    tone: str | None,
) -> None:
    """
    Apply the optional Purpose, Include, and Tone form controls before submitting.
    All interactions use JavaScript clicks for Vuetify's visually-hidden inputs.
    """
    # ── Purpose radio button ──────────────────────────────────────────────────
    if purpose:
        aria_label = PURPOSE_OPTIONS[purpose]
        radio = wait.until(EC.presence_of_element_located((
            By.CSS_SELECTOR, f'input[type="radio"][aria-label="{aria_label}"]'
        )))
        driver.execute_script("arguments[0].click();", radio)
        time.sleep(0.3)  # let Vue finish re-rendering before proceeding

    # ── Include checkboxes ────────────────────────────────────────────────────
    for key in includes:
        aria_label = INCLUDE_OPTIONS[key]
        cb = wait.until(EC.presence_of_element_located((
            By.CSS_SELECTOR, f'input[type="checkbox"][aria-label="{aria_label}"]'
        )))
        driver.execute_script("arguments[0].click();", cb)
        time.sleep(0.2)

    # ── Tone dropdown (Vuetify v-select) ──────────────────────────────────────
    # The Vuetify v-select renders the clickable trigger as the outer
    # <div role="combobox"> container, not the inner <input role="combobox">.
    # Clicking the container opens the dropdown overlay list.
    if tone:
        search_text = TONE_OPTIONS[tone].lower()   # e.g. "academic style"
        key_text    = tone.lower()                  # e.g. "academic" (fallback)

        combobox_field = wait.until(EC.element_to_be_clickable((
            By.CSS_SELECTOR, '.v-field[role="combobox"]'
        )))
        combobox_field.click()

        # Wait for the Vuetify overlay list container to appear.
        # The overlay is teleported to the document body and is scoped inside
        # .v-overlay__content — searching here avoids matching the nav sidebar.
        overlay_list = wait.until(EC.visibility_of_element_located((
            By.CSS_SELECTOR, '.v-overlay__content .v-list'
        )))
        time.sleep(0.3)  # let all list items finish rendering

        # Try .v-list-item first; fall back to [role="option"] if empty.
        items = overlay_list.find_elements(By.CSS_SELECTOR, '.v-list-item')
        if not items:
            items = overlay_list.find_elements(By.CSS_SELECTOR, '[role="option"]')

        # Use JS innerText (reflects what is actually displayed) rather than
        # the raw DOM textContent which may include hidden / whitespace nodes.
        def _item_text(el) -> str:
            t = driver.execute_script("return arguments[0].innerText", el) or ""
            return t.strip().lower()

        matched = False
        for item in items:
            text = _item_text(item)
            if search_text in text or key_text in text:
                driver.execute_script("arguments[0].click();", item)
                matched = True
                break

        if not matched:
            all_texts = ", ".join(repr(_item_text(el)) for el in items) or "(none)"
            raise RuntimeError(
                f"Could not find tone '{tone}' in the dropdown. "
                f"Found {len(items)} item(s): {all_texts}."
            )

        time.sleep(0.2)  # let Vue register the selection before submit


def upload_and_generate(
    driver: webdriver.Remote,
    image_path: Path,
    purpose: str | None = None,
    includes: list[str] | None = None,
    tone: str | None = None,
    timeout: int = GENERATION_TIMEOUT,
) -> str:
    """
    Upload one image, click Generate Content, wait for the full response,
    and return the raw page text that contains the three alt text versions.

    Raises RuntimeError if generation times out or no output is detected.
    """
    wait = WebDriverWait(driver, 30)

    # ── 1. Snapshot baseline text before any submission ───────────────────────
    # The static page already contains phrases like "a short version, medium
    # version, and long version" inside the visible system-prompt box. We must
    # record this before submitting so we can ignore it when detecting the
    # actual AI response.
    baseline: str = driver.execute_script("return document.body.innerText") or ""

    # ── 2. Upload the image ───────────────────────────────────────────────────
    # The Vuetify file-upload hides the native <input type="file"> visually.
    # Selenium can still send keys to it without making it visible.
    file_input = wait.until(
        EC.presence_of_element_located((By.CSS_SELECTOR, 'input[type="file"]'))
    )
    # Make the input interactable in case it is display:none
    driver.execute_script("arguments[0].style.display = 'block';", file_input)
    file_input.send_keys(str(image_path.resolve()))

    # Brief pause for Vue to register the file selection
    time.sleep(0.5)

    # ── 3. Apply optional form options ────────────────────────────────────────
    _set_form_options(driver, wait, purpose, includes or [], tone)

    # ── 4. Click Generate Content ─────────────────────────────────────────────
    # Use element_to_be_clickable (not just presence_of_element_located) so we
    # wait until Vue has finished processing the file upload and the button is
    # actually enabled — this matters especially when no form options are set,
    # because _set_form_options then does nothing and provides no implicit wait.
    submit_btn = wait.until(
        EC.element_to_be_clickable((By.CSS_SELECTOR, 'button[type="submit"]'))
    )
    driver.execute_script("arguments[0].scrollIntoView({block: 'center'});", submit_btn)
    time.sleep(0.2)
    driver.execute_script("arguments[0].click();", submit_btn)

    # ── 5. Wait for the AI response to appear and stabilize ───────────────────
    deadline = time.monotonic() + timeout
    last_text = ""
    stable_since = None

    while time.monotonic() < deadline:
        page_text: str = driver.execute_script("return document.body.innerText") or ""

        # Isolate the AI response text.
        # When the form collapses on submission 'startswith' fails; fall back
        # to _extract_response_text which only returns a value once the full
        # Short→Medium→Long sequence is present.  While that sequence is absent
        # (returns None) we keep polling — this prevents the loop from latching
        # onto static page content that happens to contain a version heading.
        if page_text.startswith(baseline):
            new_text = page_text[len(baseline):]
            if not _VERSION_HEADING.search(new_text):
                time.sleep(0.5)
                continue
        else:
            new_text = _extract_response_text(page_text)
            if new_text is None:
                time.sleep(0.5)
                continue

        if new_text == last_text:
            if stable_since is None:
                stable_since = time.monotonic()
            elif time.monotonic() - stable_since >= STABILITY_SECONDS:
                return new_text   # stable AI response ready
        else:
            stable_since = None  # still streaming
            last_text = new_text

        time.sleep(0.5)

    raise RuntimeError(
        f"Generation timed out after {timeout}s. "
        "No standalone 'Short/Medium/Long Version' heading appeared in new page text."
    )


# ── orchestration ─────────────────────────────────────────────────────────────

def run_batch(
    driver: webdriver.Remote,
    pptx_path: Path,
    output_path: Path,
    tool_url: str,
    version: str,
    purpose: str | None,
    includes: list[str],
    tone: str | None,
    stop_event: threading.Event | None = None,
) -> None:
    """
    Process a presentation using an already-authenticated driver.
    The caller is responsible for browser lifecycle (build_driver / driver.quit).
    """
    prs = Presentation(pptx_path)

    # Collect all picture shapes up front so we can show progress
    images: list[tuple[int, object]] = []   # (slide_num, shape)
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append((slide_num, shape))

    if not images:
        print("No picture shapes found in the presentation. Nothing to do.")
        return

    print(f"Found {len(images)} picture(s) across {len(prs.slides)} slide(s).")

    ok = 0
    skipped = 0
    errors = 0

    for idx, (slide_num, shape) in enumerate(images, start=1):
        if stop_event and stop_event.is_set():
            print("\nAborted by user.\n")
            break

        image = shape.image
        ct = image.content_type.lower()

        print(f"[{idx}/{len(images)}] Slide {slide_num} — {shape.name!r} ({ct})")

        # Save image to a temp file (converts raster and vector formats to PNG as needed)
        tmp_path = None
        try:
            tmp_path = save_image_to_temp(image.blob, ct)

            # Navigate to a fresh form (keeps the session alive)
            reset_page(driver, tool_url)

            # Upload and generate
            raw_output = upload_and_generate(
                driver, tmp_path,
                purpose=purpose, includes=includes, tone=tone,
            )

            # Extract the requested version
            alt_text = extract_version(raw_output, version)

            # Write into the PPTX XML
            wrote = set_alt_text(shape, description=alt_text)
            if wrote:
                preview = alt_text[:100] + ("…" if len(alt_text) > 100 else "")
                print(f"  OK — {preview}\n")
                ok += 1
            else:
                print("  WARNING — could not locate cNvPr element; alt text not written.\n")
                errors += 1

        except _SESSION_LOST_EXCEPTIONS:
            raise   # browser is gone — let the caller handle it
        except Exception as exc:
            print(f"  ERROR — {exc}\n")
            errors += 1

        finally:
            if tmp_path and tmp_path.exists():
                tmp_path.unlink(missing_ok=True)

    prs.save(output_path)

    print("─" * 60)
    print(f"Total images : {len(images)}")
    print(f"  Processed  : {ok}")
    print(f"  Skipped    : {skipped}")
    print(f"  Errors     : {errors}")
    print(f"Output saved : {output_path}")


def process_presentation(
    pptx_path: Path,
    output_path: Path,
    tool_url: str,
    version: str,
    browser: str,
    geckodriver_path: Path | None,
    purpose: str | None,
    includes: list[str],
    tone: str | None,
    stop_event: threading.Event | None = None,
) -> None:
    driver = build_driver(browser, geckodriver_path)
    try:
        wait_for_auth(driver, tool_url)
        run_batch(
            driver=driver,
            pptx_path=pptx_path,
            output_path=output_path,
            tool_url=tool_url,
            version=version,
            purpose=purpose,
            includes=includes,
            tone=tone,
            stop_event=stop_event,
        )
    finally:
        driver.quit()


# ── entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Batch-generate alt text for every picture in a .pptx file "
            "using the UMich AI Helper Alt Text Generator (browser-based, "
            "institutional Google authentication)."
        )
    )
    parser.add_argument("pptx", type=Path, help="Input .pptx file")
    parser.add_argument(
        "--url", "-u", metavar="URL",
        default="https://aihelper.engin.umich.edu/alt-text-generator",
        help=(
            "Full URL of the Alt Text Generator page. "
            "(default: https://aihelper.engin.umich.edu/alt-text-generator)"
        ),
    )
    parser.add_argument(
        "--output", "-o", type=Path, default=None,
        help="Output .pptx path (default: <stem>_alt_text.pptx)",
    )
    parser.add_argument(
        "--version", "-v",
        choices=["short", "medium", "long"],
        default="long",
        help=(
            "Which of the three alt text lengths to embed: "
            "short ≈ 1 sentence, medium ≈ 2-3 sentences, long = full description. "
            "(default: long)"
        ),
    )
    parser.add_argument(
        "--browser", "-b",
        choices=["auto", "firefox", "chrome", "edge"],
        default="auto",
        help=(
            "Browser to use. 'auto' tries Firefox → Chrome → Edge. "
            "(default: auto)"
        ),
    )
    parser.add_argument(
        "--geckodriver", "-g", type=Path, default=None,
        help=(
            "Path to geckodriver executable (Firefox only). "
            "Defaults to geckodriver.exe in the script's folder, "
            "then falls back to system PATH."
        ),
    )
    parser.add_argument(
        "--purpose", "-p",
        choices=list(PURPOSE_OPTIONS),
        default=None,
        metavar="{" + ",".join(PURPOSE_OPTIONS) + "}",
        help=(
            "Purpose/Use Case radio button to select on the form. "
            "Choices: general, educational, instructional, marketing, icon. "
            "(default: none selected)"
        ),
    )
    parser.add_argument(
        "--include", "-i",
        choices=list(INCLUDE_OPTIONS),
        nargs="*",
        default=[],
        dest="includes",
        metavar="{" + ",".join(INCLUDE_OPTIONS) + "}",
        help=(
            "Optional checkboxes to tick. "
            "Choices: data-values, captions. "
            "Can be repeated or space-separated. "
            "(default: none selected)"
        ),
    )
    parser.add_argument(
        "--tone", "-t",
        choices=list(TONE_OPTIONS),
        default=None,
        metavar="{" + ",".join(TONE_OPTIONS) + "}",
        help=(
            "Tone of the generated alt text. "
            "Choices: formal, academic, professional, neutral, "
            "conversational, casual, colloquial. "
            "(default: none selected)"
        ),
    )
    args = parser.parse_args()

    if not args.pptx.exists():
        sys.exit(f"File not found: {args.pptx}")

    output_path = args.output or args.pptx.with_stem(args.pptx.stem + "_alt_text")

    # Resolve geckodriver location
    geckodriver = args.geckodriver
    if geckodriver is None and _LOCAL_GECKODRIVER.exists():
        geckodriver = _LOCAL_GECKODRIVER

    print(f"Input    : {args.pptx}")
    print(f"Output   : {output_path}")
    print(f"Version  : {args.version}")
    print(f"Purpose  : {args.purpose or '(none)'}")
    print(f"Include  : {', '.join(args.includes) if args.includes else '(none)'}")
    print(f"Tone     : {args.tone or '(none)'}")
    print(f"Tool URL : {args.url}")
    if geckodriver:
        print(f"Driver   : {geckodriver}")

    process_presentation(
        pptx_path=args.pptx,
        output_path=output_path,
        tool_url=args.url,
        version=args.version,
        browser=args.browser,
        geckodriver_path=geckodriver,
        purpose=args.purpose,
        includes=args.includes,
        tone=args.tone,
    )


if __name__ == "__main__":
    main()
