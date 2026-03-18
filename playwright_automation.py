#!/usr/bin/env python3
"""
playwright_automation.py  –  Playwright-based automation for the AI Helper tool.

Used by web_app.py instead of the Selenium automation in generate_alt_text.py.
Key differences from the Selenium approach:

  • Per-user sessions: each user's Playwright storageState (cookies + localStorage)
    is saved to SESSIONS_DIR/<hash>.json after their first sign-in.
  • Persistent across restarts: when SESSIONS_DIR is a PersistentVolumeClaim the
    session survives container restarts; users only re-authenticate when Google
    actually expires their session.
  • Headless processing: after the one-time sign-in, all processing runs in a
    headless Chromium instance — no display or VNC needed.
  • Interactive sign-in: the initial sign-in launches a non-headless Chromium on
    the Xvfb display (:99) so the user can complete Google OAuth through the
    noVNC window.  Only one sign-in can proceed at a time (serialised with a lock).
"""

import hashlib
import json
import os
import sys
import threading
import time
from pathlib import Path

try:
    from playwright.sync_api import sync_playwright, TimeoutError as PWTimeout
except ImportError:
    import sys
    sys.exit(
        "Missing dependency: pip install playwright && playwright install chromium"
    )

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Import the pure-Python helpers from the CLI module.
# (generate_alt_text.py also imports Selenium at the top level, but Selenium is
# still in requirements.txt for CLI use, so the import succeeds in the container.)
from generate_alt_text import (
    PURPOSE_OPTIONS,
    INCLUDE_OPTIONS,
    TONE_OPTIONS,
    save_image_to_temp,
    extract_version,
    set_alt_text,
    _VERSION_HEADING,
    _extract_response_text,
)

# ── Configuration ──────────────────────────────────────────────────────────────

SESSIONS_DIR       = Path(os.environ.get("SESSIONS_DIR") or "./sessions")
AUTH_TIMEOUT_MS    = 300_000   # 5 minutes — Playwright uses milliseconds
GENERATION_TIMEOUT = 60        # seconds per image
STABILITY_SECONDS  = 1.0       # seconds of unchanged text before response is "done"

# Required in Docker/OpenShift containers; harmless locally.
_CHROMIUM_ARGS = ["--no-sandbox", "--disable-dev-shm-usage"]


def _launch_kwargs() -> dict:
    """
    Return kwargs for p.chromium.launch().

    Both sign-in and processing run non-headless to pass Cloudflare's bot
    detection (which blocks headless Chromium regardless of session cookies).
    On Linux the Xvfb virtual display (:99) is used so the window stays
    invisible to users; on Windows it opens on the desktop (fine for testing).
    """
    kwargs: dict = {"headless": False, "args": _CHROMIUM_ARGS}
    if sys.platform != "win32":
        kwargs["env"] = {**os.environ, "DISPLAY": ":99"}
    return kwargs


# ── Session file helpers ───────────────────────────────────────────────────────

def _session_path(user_id: str) -> Path:
    h = hashlib.sha256(user_id.encode()).hexdigest()[:24]
    return SESSIONS_DIR / f"{h}.json"


def has_session(user_id: str) -> bool:
    return _session_path(user_id).exists()


def delete_session(user_id: str) -> None:
    p = _session_path(user_id)
    if p.exists():
        p.unlink()


def _save_storage_state(context, user_id: str) -> None:
    SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
    _session_path(user_id).write_text(json.dumps(context.storage_state()))


# ── Interactive sign-in (one user at a time) ───────────────────────────────────

_signin_lock   = threading.Lock()
_signin_owner: str | None = None         # user_id currently holding the lock
_signin_status: dict[str, dict] = {}     # {user_id: {"status": ..., "msg": ...}}
_status_guard  = threading.Lock()        # protects _signin_owner and _signin_status


def set_signin_status(user_id: str, status: str, msg: str = "") -> None:
    with _status_guard:
        _signin_status[user_id] = {"status": status, "msg": msg}


def get_signin_status(user_id: str) -> dict:
    with _status_guard:
        return dict(_signin_status.get(user_id, {"status": "none", "msg": ""}))


def signin_lock_available() -> bool:
    return not _signin_lock.locked()


def start_signin(user_id: str, tool_url: str) -> bool:
    """
    Try to start an interactive browser sign-in for *user_id*.
    Returns True if the background thread was launched.
    Returns False if another user is currently signing in (caller should ask the
    user to try again in a few minutes).
    The browser opens on DISPLAY=:99 (Xvfb) and is visible via noVNC on port 6080.
    """
    global _signin_owner

    if not _signin_lock.acquire(blocking=False):
        return False

    with _status_guard:
        _signin_owner = user_id
    set_signin_status(user_id, "launching")

    def _worker() -> None:
        global _signin_owner
        try:
            with sync_playwright() as p:
                browser = p.chromium.launch(**_launch_kwargs())
                context = browser.new_context()
                page    = context.new_page()
                page.goto(tool_url)
                set_signin_status(
                    user_id, "waiting",
                    "Please sign in with your UMich Google account "
                    "in the browser window.",
                )
                try:
                    page.wait_for_selector(
                        'a[href="/profile"]', timeout=AUTH_TIMEOUT_MS
                    )
                    _save_storage_state(context, user_id)
                    set_signin_status(
                        user_id, "ready", "Signed in ✓  Ready to process files."
                    )
                except PWTimeout:
                    set_signin_status(
                        user_id, "timeout",
                        "Sign-in timed out. Click 'Connect' to try again.",
                    )
                finally:
                    browser.close()
        except Exception as exc:
            set_signin_status(user_id, "error", str(exc))
        finally:
            with _status_guard:
                _signin_owner = None
            try:
                _signin_lock.release()
            except RuntimeError:
                pass

    threading.Thread(target=_worker, daemon=True).start()
    return True


# ── Headless batch processing ──────────────────────────────────────────────────

class AuthExpiredError(RuntimeError):
    """The stored session is missing or no longer valid."""


def run_batch_headless(
    user_id:     str,
    pptx_path:   Path,
    output_path: Path,
    tool_url:    str,
    version:     str,
    purpose:     str | None,
    includes:    list[str],
    tone:        str | None,
    stop_event:  threading.Event | None = None,
    log_fn=None,
) -> None:
    """
    Process a presentation using the user's stored Playwright session state.
    Raises AuthExpiredError when the session is missing or has expired.
    """
    def log(msg: str) -> None:
        (log_fn or print)(msg)

    if not has_session(user_id):
        raise AuthExpiredError(
            "No session found. Please sign in to the AI Helper tool first."
        )

    try:
        state = json.loads(_session_path(user_id).read_text())
    except Exception as exc:
        raise AuthExpiredError(f"Could not load session: {exc}")

    with sync_playwright() as p:
        browser = p.chromium.launch(**_launch_kwargs())
        context = browser.new_context(storage_state=state)
        page    = context.new_page()

        # Verify the session is still valid before doing any real work.
        # wait_for_load_state("networkidle") lets the Vue SPA finish all its
        # requests before we look for the profile link, avoiding false negatives
        # from a race between page initialisation and the timeout clock.
        page.goto(tool_url, timeout=30_000)
        try:
            page.wait_for_load_state("networkidle", timeout=20_000)
        except PWTimeout:
            pass  # proceed anyway; the selector check below is the real gate
        try:
            page.wait_for_selector('a[href="/profile"]', timeout=20_000)
        except PWTimeout:
            browser.close()
            delete_session(user_id)
            raise AuthExpiredError("Session expired. Please sign in again.")

        prs    = Presentation(pptx_path)
        images = [
            (slide_num, shape)
            for slide_num, slide in enumerate(prs.slides, start=1)
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]

        if not images:
            log("No picture shapes found in the presentation. Nothing to do.")
            browser.close()
            prs.save(output_path)
            return

        log(f"Found {len(images)} picture(s) across {len(prs.slides)} slide(s).")

        ok = errors = 0

        for idx, (slide_num, shape) in enumerate(images, start=1):
            if stop_event and stop_event.is_set():
                log("\nAborted by user.\n")
                break

            image = shape.image
            ct    = image.content_type.lower()
            log(f"[{idx}/{len(images)}] Slide {slide_num} — {shape.name!r} ({ct})")

            tmp_path = None
            try:
                tmp_path = save_image_to_temp(image.blob, ct)

                # Fresh form for each image.
                page.goto(tool_url, timeout=30_000)
                try:
                    page.wait_for_load_state("networkidle", timeout=20_000)
                except PWTimeout:
                    pass
                try:
                    page.wait_for_selector('input[type="file"]', timeout=30_000)
                except PWTimeout:
                    raise RuntimeError("Timed out waiting for the form to load.")

                raw      = _drive_form(page, tmp_path, purpose, includes, tone)
                alt_text = extract_version(raw, version)
                wrote    = set_alt_text(shape, description=alt_text)

                if wrote:
                    preview = alt_text[:100] + ("…" if len(alt_text) > 100 else "")
                    log(f"  OK — {preview}\n")
                    ok += 1
                else:
                    log("  WARNING — could not locate cNvPr element; "
                        "alt text not written.\n")
                    errors += 1

            except Exception as exc:
                log(f"  ERROR — {exc}\n")
                errors += 1
            finally:
                if tmp_path and tmp_path.exists():
                    tmp_path.unlink(missing_ok=True)

        prs.save(output_path)
        browser.close()

    log("─" * 60)
    log(f"Total images : {len(images)}")
    log(f"  Processed  : {ok}")
    log(f"  Errors     : {errors}")
    log(f"Output saved : {output_path}")


def _drive_form(
    page,
    image_path: Path,
    purpose:    str | None,
    includes:   list[str],
    tone:       str | None,
    timeout:    int = GENERATION_TIMEOUT,
) -> str:
    """Drive the AI Helper web form for one image and return the raw response text."""
    baseline: str = page.evaluate("document.body.innerText") or ""

    # Upload
    fi = page.locator('input[type="file"]')
    fi.evaluate("el => el.style.display = 'block'")
    fi.set_input_files(str(image_path.resolve()))
    time.sleep(0.5)

    # Purpose radio
    if purpose:
        page.locator(
            f'input[type="radio"][aria-label="{PURPOSE_OPTIONS[purpose]}"]'
        ).evaluate("el => el.click()")
        time.sleep(0.3)

    # Include checkboxes
    for key in (includes or []):
        page.locator(
            f'input[type="checkbox"][aria-label="{INCLUDE_OPTIONS[key]}"]'
        ).evaluate("el => el.click()")
        time.sleep(0.2)

    # Tone dropdown
    if tone:
        search = TONE_OPTIONS[tone].lower()
        page.locator('.v-field[role="combobox"]').click()
        overlay = page.locator('.v-overlay__content .v-list')
        overlay.wait_for(state="visible", timeout=10_000)
        time.sleep(0.3)
        items = overlay.locator('.v-list-item').all()
        if not items:
            items = overlay.locator('[role="option"]').all()
        matched = False
        for item in items:
            text = (item.evaluate("el => el.innerText") or "").strip().lower()
            if search in text or tone.lower() in text:
                item.evaluate("el => el.click()")
                matched = True
                break
        if not matched:
            texts = ", ".join(
                repr((i.evaluate("el => el.innerText") or "").strip())
                for i in items
            ) or "(none)"
            raise RuntimeError(
                f"Tone '{tone}' not found in dropdown. Got: {texts}"
            )
        time.sleep(0.2)

    # Submit
    btn = page.locator('button[type="submit"]')
    btn.wait_for(state="visible", timeout=30_000)
    btn.evaluate("el => { el.scrollIntoView({block:'center'}); el.click(); }")
    time.sleep(0.2)

    # Wait for stable AI response
    deadline     = time.monotonic() + timeout
    last_text    = ""
    stable_since: float | None = None

    while time.monotonic() < deadline:
        page_text: str = page.evaluate("document.body.innerText") or ""

        if page_text.startswith(baseline):
            new_text = page_text[len(baseline):]
            if not _VERSION_HEADING.search(new_text):
                time.sleep(0.5)
                continue
        else:
            new_text = _extract_response_text(page_text)
            if new_text is None:
                time.sleep(0.5)
                continue

        if new_text == last_text:
            if stable_since is None:
                stable_since = time.monotonic()
            elif time.monotonic() - stable_since >= STABILITY_SECONDS:
                return new_text
        else:
            stable_since = None
            last_text    = new_text

        time.sleep(0.5)

    raise RuntimeError(
        f"Generation timed out after {timeout}s. "
        "No Short/Medium/Long heading sequence appeared in new page text."
    )
